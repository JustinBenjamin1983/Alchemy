# DDCheckReadability/__init__.py
"""
Document Readability Check Endpoint

Validates that documents can be read/parsed before the main DD processing begins.
This is a lightweight check that verifies:
1. File exists in blob storage (or local storage in DEV_MODE)
2. File is not empty
3. File extension is supported
4. File can be opened (not corrupted or password-protected)
"""
import logging
import os
import json
import io
import azure.functions as func

import uuid

from shared.utils import auth_get_email
from shared.session import transactional_session
from shared.models import Document, Folder, DueDiligence
from shared.uploader import read_from_blob_storage, write_to_blob_storage
from sqlalchemy.orm import joinedload

from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.units import inch

DEV_MODE = os.environ.get("DEV_MODE", "").lower() == "true"
LOCAL_STORAGE_PATH = os.environ.get("LOCAL_STORAGE_PATH", "")

# Supported file types for DD processing
SUPPORTED_TYPES = {"pdf", "docx", "pptx", "xlsx", "jpg", "jpeg", "png", "bmp", "tiff", "doc", "xls", "ppt"}


def convert_pptx_to_pdf(file_contents: bytes, filename: str) -> bytes:
    """
    Extract text from PPTX and create a readable PDF.
    Returns PDF bytes.
    """
    from pptx import Presentation

    pptx_stream = io.BytesIO(file_contents)
    prs = Presentation(pptx_stream)

    # Create PDF
    pdf_buffer = io.BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # Add slide header
        story.append(Paragraph(f"<b>Slide {slide_num}</b>", styles['Heading2']))
        story.append(Spacer(1, 12))

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Escape special XML characters for reportlab
                text = shape.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                try:
                    story.append(Paragraph(text, styles['Normal']))
                    story.append(Spacer(1, 6))
                except Exception as e:
                    # If paragraph fails, try plain text without formatting
                    logging.warning(f"Failed to add formatted text, using plain: {e}")
                    story.append(Paragraph(text[:500], styles['Normal']))  # Truncate if needed

        story.append(Spacer(1, 20))

    doc.build(story)
    return pdf_buffer.getvalue()


def convert_docx_to_pdf(file_contents: bytes, filename: str) -> bytes:
    """
    Extract text from DOCX and create a readable PDF.
    Preserves basic formatting (bold, italic, underline) and paragraph structure.
    Returns PDF bytes.
    """
    from docx import Document as DocxDocument
    from docx.shared import Pt
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

    docx_stream = io.BytesIO(file_contents)
    doc = DocxDocument(docx_stream)

    # Create PDF
    pdf_buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        pdf_buffer,
        pagesize=letter,
        leftMargin=0.75*inch,
        rightMargin=0.75*inch,
        topMargin=0.75*inch,
        bottomMargin=0.75*inch
    )
    styles = getSampleStyleSheet()

    # Create custom styles
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading1'],
        fontSize=14,
        spaceAfter=12
    )
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=11,
        leading=14,
        spaceAfter=6
    )

    story = []

    for para in doc.paragraphs:
        if not para.text.strip():
            story.append(Spacer(1, 6))
            continue

        # Build formatted text with inline styles
        formatted_text = ""
        for run in para.runs:
            text = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            if not text:
                continue

            # Apply formatting tags
            if run.bold and run.italic:
                text = f"<b><i>{text}</i></b>"
            elif run.bold:
                text = f"<b>{text}</b>"
            elif run.italic:
                text = f"<i>{text}</i>"
            if run.underline:
                text = f"<u>{text}</u>"

            formatted_text += text

        if not formatted_text.strip():
            continue

        # Determine if this is a heading based on style name
        try:
            style_name = para.style.name.lower() if para.style else ""
            if 'heading' in style_name:
                story.append(Paragraph(formatted_text, heading_style))
            else:
                story.append(Paragraph(formatted_text, normal_style))
        except Exception as e:
            logging.warning(f"Failed to add paragraph, using plain text: {e}")
            # Fallback: plain text without formatting
            plain_text = para.text[:1000].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(plain_text, normal_style))

    # Handle tables
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text[:200].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                row_data.append(Paragraph(cell_text, styles['Normal']))
            table_data.append(row_data)

        if table_data:
            try:
                # Calculate column widths
                num_cols = max(len(row) for row in table_data) if table_data else 1
                col_width = (7 * inch) / num_cols

                pdf_table = Table(table_data, colWidths=[col_width] * num_cols)
                pdf_table.setStyle(TableStyle([
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('PADDING', (0, 0), (-1, -1), 4),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                story.append(Spacer(1, 12))
                story.append(pdf_table)
                story.append(Spacer(1, 12))
            except Exception as e:
                logging.warning(f"Failed to add table: {e}")

    if not story:
        story.append(Paragraph("(Document appears to be empty)", styles['Normal']))

    pdf_doc.build(story)
    return pdf_buffer.getvalue()


def convert_xlsx_to_pdf(file_contents: bytes, filename: str) -> bytes:
    """
    Extract data from XLSX and create a readable PDF with tables.
    Each sheet becomes a section in the PDF.
    Returns PDF bytes.
    """
    from openpyxl import load_workbook

    xlsx_stream = io.BytesIO(file_contents)
    wb = load_workbook(xlsx_stream, read_only=True, data_only=True)

    # Create PDF in landscape orientation for better table display
    pdf_buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(
        pdf_buffer,
        pagesize=landscape(letter),
        leftMargin=0.5*inch,
        rightMargin=0.5*inch,
        topMargin=0.5*inch,
        bottomMargin=0.5*inch
    )
    styles = getSampleStyleSheet()
    story = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Add sheet header
        story.append(Paragraph(f"<b>Sheet: {sheet_name}</b>", styles['Heading2']))
        story.append(Spacer(1, 12))

        # Collect data from sheet (limit to reasonable size)
        table_data = []
        max_rows = 500  # Limit rows per sheet
        max_cols = 20   # Limit columns

        for row_idx, row in enumerate(sheet.iter_rows(max_row=max_rows, max_col=max_cols)):
            if row_idx > max_rows:
                break

            row_data = []
            has_data = False
            for cell in row:
                value = cell.value
                if value is not None:
                    has_data = True
                    # Format cell value
                    if isinstance(value, (int, float)):
                        text = str(value)
                    else:
                        text = str(value)[:100]  # Truncate long text
                    text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                else:
                    text = ""
                row_data.append(text)

            # Only add rows that have data
            if has_data:
                table_data.append(row_data)

        if table_data:
            try:
                num_cols = max(len(row) for row in table_data)

                # Limit columns to fit page width - max ~12 columns in landscape
                max_cols_per_page = 12
                if num_cols > max_cols_per_page:
                    num_cols = max_cols_per_page

                # Fixed column width to fill page (10" / num_cols)
                col_width = (10 * inch) / num_cols

                # Create paragraph style that wraps text within cell width
                cell_style = ParagraphStyle(
                    'CellStyle',
                    parent=styles['Normal'],
                    fontSize=8,
                    leading=10,
                    wordWrap='CJK',  # Enables word wrapping
                )

                # Format data with Paragraph objects for text wrapping
                formatted_data = []
                for row in table_data[:300]:  # Limit rows
                    formatted_row = []
                    for col_idx, cell in enumerate(row[:num_cols]):
                        text = str(cell) if cell else ""
                        # Escape special characters for reportlab
                        text = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        formatted_row.append(Paragraph(text, cell_style))
                    # Pad row to have consistent columns
                    while len(formatted_row) < num_cols:
                        formatted_row.append(Paragraph("", cell_style))
                    formatted_data.append(formatted_row)

                pdf_table = Table(formatted_data, colWidths=[col_width] * num_cols)
                pdf_table.setStyle(TableStyle([
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                    ('LEFTPADDING', (0, 0), (-1, -1), 3),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 3),
                    ('TOPPADDING', (0, 0), (-1, -1), 2),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                story.append(pdf_table)

                if len(table_data) > 200:
                    story.append(Spacer(1, 6))
                    story.append(Paragraph(f"<i>(Showing first 200 of {len(table_data)} rows)</i>", styles['Normal']))

            except Exception as e:
                logging.warning(f"Failed to create table for sheet {sheet_name}: {e}")
                story.append(Paragraph(f"(Error rendering sheet: {str(e)[:100]})", styles['Normal']))
        else:
            story.append(Paragraph("(Sheet is empty)", styles['Normal']))

        story.append(PageBreak())

    if not story:
        story.append(Paragraph("(Workbook appears to be empty)", styles['Normal']))

    pdf_doc.build(story)
    return pdf_buffer.getvalue()


def store_converted_pdf(original_doc_id: str, pdf_bytes: bytes, original_filename: str, folder_id: str, session) -> str:
    """
    Store converted PDF and create Document record.
    Returns the new document ID.
    """
    converted_doc_id = str(uuid.uuid4())
    pdf_filename = original_filename.rsplit('.', 1)[0] + '_converted.pdf'

    # Store PDF (local in DEV_MODE, blob storage otherwise)
    if DEV_MODE and LOCAL_STORAGE_PATH:
        local_path = os.path.join(LOCAL_STORAGE_PATH, "docs", converted_doc_id)
        os.makedirs(os.path.dirname(local_path), exist_ok=True)
        with open(local_path, "wb") as f:
            f.write(pdf_bytes)
        logging.info(f"[DEV MODE] Stored converted PDF at: {local_path}")
    else:
        write_to_blob_storage(
            os.environ["DD_DOCS_BLOB_STORAGE_CONNECTION_STRING"],
            os.environ.get("DD_DOCS_STORAGE_CONTAINER_NAME", "dd-docs"),
            converted_doc_id,
            pdf_bytes,
            meta_data={
                "original_file_name": pdf_filename,
                "extension": "pdf",
                "converted_from": original_doc_id,
                "is_converted": "true"
            }
        )
        logging.info(f"Stored converted PDF in blob storage: {converted_doc_id}")

    # Create new Document record for the converted PDF
    new_doc = Document(
        id=uuid.UUID(converted_doc_id),
        folder_id=uuid.UUID(folder_id) if isinstance(folder_id, str) else folder_id,
        type="pdf",
        original_file_name=pdf_filename,
        processing_status="pending",
        is_original=False,
        size_in_bytes=len(pdf_bytes),
        readability_status="ready",  # Converted PDF is inherently readable
        converted_from_id=uuid.UUID(original_doc_id) if isinstance(original_doc_id, str) else original_doc_id
    )
    session.add(new_doc)
    session.flush()  # Ensure ID is available

    return converted_doc_id


def read_file_contents(doc_id: str) -> bytes:
    """
    Read file contents from storage with fallback.
    Tries primary storage first, then falls back to the other.
    This ensures files work regardless of where they were uploaded.
    """
    local_path = os.path.join(LOCAL_STORAGE_PATH, "docs", doc_id) if LOCAL_STORAGE_PATH else None
    local_path_alt = os.path.join(LOCAL_STORAGE_PATH, doc_id) if LOCAL_STORAGE_PATH else None

    def try_local_storage():
        """Attempt to read from local storage."""
        if not LOCAL_STORAGE_PATH:
            return None
        if local_path and os.path.exists(local_path):
            logging.info(f"[LOCAL] Reading document from: {local_path}")
            with open(local_path, "rb") as f:
                return f.read()
        if local_path_alt and os.path.exists(local_path_alt):
            logging.info(f"[LOCAL] Reading document from: {local_path_alt}")
            with open(local_path_alt, "rb") as f:
                return f.read()
        return None

    def try_blob_storage():
        """Attempt to read from Azure Blob Storage."""
        try:
            conn_string = os.environ.get("DD_DOCS_BLOB_STORAGE_CONNECTION_STRING")
            if not conn_string:
                logging.warning("[BLOB] No blob storage connection string configured")
                return None
            logging.info(f"[BLOB] Reading document from Azure Blob Storage: {doc_id}")
            return read_from_blob_storage(
                conn_string,
                os.environ.get("DD_DOCS_STORAGE_CONTAINER_NAME", "dd-docs"),
                doc_id
            )
        except Exception as e:
            logging.warning(f"[BLOB] Failed to read from blob storage: {e}")
            return None

    # Try primary storage first, then fallback
    if DEV_MODE:
        # DEV_MODE: Try local first, then blob
        content = try_local_storage()
        if content is not None:
            return content
        logging.info(f"[FALLBACK] File not in local storage, trying Azure Blob...")
        content = try_blob_storage()
        if content is not None:
            return content
        raise FileNotFoundError(f"File not found in local storage or blob storage: {doc_id}")
    else:
        # Production: Try blob first, then local
        content = try_blob_storage()
        if content is not None:
            return content
        logging.info(f"[FALLBACK] File not in blob storage, trying local...")
        content = try_local_storage()
        if content is not None:
            return content
        raise FileNotFoundError(f"File not found in blob storage or local storage: {doc_id}")


def check_pdf_readability(file_contents: bytes, filename: str) -> tuple[bool, str]:
    """Check if a PDF file is readable (not corrupted or password-protected)."""
    try:
        import fitz  # PyMuPDF

        pdf_stream = io.BytesIO(file_contents)
        doc = fitz.open(stream=pdf_stream, filetype="pdf")

        # Check if encrypted/password-protected
        if doc.is_encrypted:
            doc.close()
            return False, "Document is password-protected and cannot be accessed"

        # Try to read first page to verify it's not corrupted
        if doc.page_count == 0:
            doc.close()
            return False, "Document contains no pages"

        # Try to extract text from first page
        try:
            first_page = doc[0]
            _ = first_page.get_text()
        except Exception as e:
            doc.close()
            return False, f"Unable to read document content: {str(e)}"

        doc.close()
        return True, ""

    except Exception as e:
        return False, f"Unable to open document: {str(e)}"


def check_docx_readability(
    file_contents: bytes,
    filename: str,
    doc_id: str = None,
    folder_id: str = None,
    session = None
) -> tuple[bool, str, str | None]:
    """
    Check if a DOCX file is readable and convert to PDF.
    Returns (is_readable, error_message, converted_doc_id).
    """
    try:
        from docx import Document as DocxDocument

        docx_stream = io.BytesIO(file_contents)
        doc = DocxDocument(docx_stream)

        # Try to read paragraphs
        paragraph_count = len(doc.paragraphs)
        if paragraph_count == 0 and len(doc.tables) == 0:
            # Check if document has any content at all
            text_content = ""
            for para in doc.paragraphs:
                text_content += para.text
            if not text_content.strip():
                return True, "", None  # Empty but valid document

        # If we have session info, perform conversion
        converted_doc_id = None
        if doc_id and folder_id and session:
            try:
                logging.info(f"Converting DOCX to PDF: {filename}")
                pdf_bytes = convert_docx_to_pdf(file_contents, filename)
                converted_doc_id = store_converted_pdf(
                    doc_id, pdf_bytes, filename, folder_id, session
                )
                logging.info(f"Successfully converted {filename} to PDF: {converted_doc_id}")
            except Exception as conv_error:
                logging.error(f"Failed to convert DOCX to PDF: {conv_error}")
                # Conversion failure is not a readability failure

        return True, "", converted_doc_id

    except Exception as e:
        error_str = str(e).lower()
        if "encrypted" in error_str or "password" in error_str:
            return False, "Document is password-protected and cannot be accessed", None
        return False, f"Unable to open document: {str(e)}", None


def check_xlsx_readability(
    file_contents: bytes,
    filename: str,
    doc_id: str = None,
    folder_id: str = None,
    session = None
) -> tuple[bool, str, str | None]:
    """
    Check if an XLSX file is readable and convert to PDF.
    Returns (is_readable, error_message, converted_doc_id).
    """
    try:
        import openpyxl

        xlsx_stream = io.BytesIO(file_contents)
        workbook = openpyxl.load_workbook(xlsx_stream, read_only=True)

        # Check if there are any sheets
        if len(workbook.sheetnames) == 0:
            workbook.close()
            return False, "Spreadsheet contains no worksheets", None

        workbook.close()

        # If we have session info, perform conversion
        converted_doc_id = None
        if doc_id and folder_id and session:
            try:
                logging.info(f"Converting XLSX to PDF: {filename}")
                pdf_bytes = convert_xlsx_to_pdf(file_contents, filename)
                converted_doc_id = store_converted_pdf(
                    doc_id, pdf_bytes, filename, folder_id, session
                )
                logging.info(f"Successfully converted {filename} to PDF: {converted_doc_id}")
            except Exception as conv_error:
                logging.error(f"Failed to convert XLSX to PDF: {conv_error}")
                # Conversion failure is not a readability failure

        return True, "", converted_doc_id

    except Exception as e:
        error_str = str(e).lower()
        if "encrypted" in error_str or "password" in error_str:
            return False, "Spreadsheet is password-protected and cannot be accessed", None
        return False, f"Unable to open spreadsheet: {str(e)}", None


def check_pptx_readability(
    file_contents: bytes,
    filename: str,
    doc_id: str = None,
    folder_id: str = None,
    session = None
) -> tuple[bool, str, str | None]:
    """
    Check if a PPTX file is readable and convert to PDF.
    Returns (is_readable, error_message, converted_doc_id).
    """
    try:
        from pptx import Presentation

        pptx_stream = io.BytesIO(file_contents)
        prs = Presentation(pptx_stream)

        # Check if there are any slides
        if len(prs.slides) == 0:
            return False, "Presentation contains no slides", None

        # If we have session info, perform conversion
        converted_doc_id = None
        if doc_id and folder_id and session:
            try:
                logging.info(f"Converting PPTX to PDF: {filename}")
                pdf_bytes = convert_pptx_to_pdf(file_contents, filename)
                converted_doc_id = store_converted_pdf(
                    doc_id, pdf_bytes, filename, folder_id, session
                )
                logging.info(f"Successfully converted {filename} to PDF: {converted_doc_id}")
            except Exception as conv_error:
                logging.error(f"Failed to convert PPTX to PDF: {conv_error}")
                # Conversion failure is not a readability failure
                # The PPTX is still readable, just couldn't convert

        return True, "", converted_doc_id

    except Exception as e:
        error_str = str(e).lower()
        if "encrypted" in error_str or "password" in error_str:
            return False, "Presentation is password-protected and cannot be accessed", None
        return False, f"Unable to open presentation: {str(e)}", None


def check_image_readability(file_contents: bytes, filename: str) -> tuple[bool, str]:
    """Check if an image file is readable."""
    try:
        from PIL import Image

        img_stream = io.BytesIO(file_contents)
        img = Image.open(img_stream)
        img.verify()  # Verify it's a valid image

        return True, ""

    except Exception as e:
        return False, f"Unable to open image: {str(e)}"


def check_document_readability(
    doc_id: str,
    file_type: str,
    filename: str,
    folder_id: str = None,
    session = None
) -> tuple[bool, str, str | None]:
    """
    Check if a document can be read.
    Returns (is_readable, error_message, converted_doc_id).
    converted_doc_id is set only for PPTX files that were converted to PDF.
    """
    try:
        # Read file from storage (local in DEV_MODE, blob storage otherwise)
        file_contents = read_file_contents(doc_id)

        if not file_contents:
            return False, "Document file is empty or could not be retrieved", None

        if len(file_contents) < 100:  # Suspiciously small file
            return False, "Document file appears to be corrupted (file too small)", None

        file_type_lower = file_type.lower()

        # Check if supported
        if file_type_lower not in SUPPORTED_TYPES:
            return False, f"Document type '{file_type}' is not supported for analysis", None

        # Type-specific checks
        if file_type_lower == "pdf":
            is_readable, error = check_pdf_readability(file_contents, filename)
            return is_readable, error, None
        elif file_type_lower in ("docx", "doc"):
            if file_type_lower == "doc":
                # .doc files need conversion, treat as potentially readable
                return True, "", None
            # DOCX files get converted to PDF
            return check_docx_readability(
                file_contents, filename, doc_id, folder_id, session
            )
        elif file_type_lower in ("xlsx", "xls"):
            if file_type_lower == "xls":
                return True, "", None  # Legacy format, assume readable
            # XLSX files get converted to PDF
            return check_xlsx_readability(
                file_contents, filename, doc_id, folder_id, session
            )
        elif file_type_lower in ("pptx", "ppt"):
            if file_type_lower == "ppt":
                return True, "", None  # Legacy format, assume readable
            # PPTX files get converted to PDF
            return check_pptx_readability(
                file_contents, filename, doc_id, folder_id, session
            )
        elif file_type_lower in ("jpg", "jpeg", "png", "bmp", "tiff"):
            is_readable, error = check_image_readability(file_contents, filename)
            return is_readable, error, None
        else:
            # Unknown but potentially supported type
            return True, "", None

    except Exception as e:
        logging.error(f"Error checking readability for {doc_id}: {str(e)}")
        return False, f"Error checking document: {str(e)}", None


def main(req: func.HttpRequest) -> func.HttpResponse:
    """
    Check readability of documents in a DD project.

    POST /api/dd-check-readability
    Body: {
        "dd_id": "uuid",
        "doc_ids": ["uuid1", "uuid2"]  // Optional - if not provided, checks all docs
    }

    Returns: {
        "dd_id": "uuid",
        "total_documents": 10,
        "results": [
            {
                "doc_id": "uuid",
                "filename": "example.pdf",
                "status": "ready|failed|checking",
                "error": "Error message if failed"
            }
        ],
        "summary": {
            "ready": 8,
            "failed": 2,
            "checking": 0
        }
    }
    """

    # Auth check
    if not DEV_MODE and req.headers.get('function-key') != os.environ.get("FUNCTION_KEY"):
        logging.info("No matching value in header")
        return func.HttpResponse("", status_code=401)

    try:
        email, err = auth_get_email(req)
        if err:
            return err

        # Parse request body
        try:
            req_body = req.get_json()
        except ValueError:
            return func.HttpResponse(
                json.dumps({"error": "Invalid JSON body"}),
                status_code=400,
                mimetype="application/json"
            )

        dd_id = req_body.get("dd_id")
        doc_ids = req_body.get("doc_ids")  # Optional list of specific doc IDs

        if not dd_id:
            return func.HttpResponse(
                json.dumps({"error": "dd_id is required"}),
                status_code=400,
                mimetype="application/json"
            )

        results = []
        summary = {"ready": 0, "failed": 0, "checking": 0}

        with transactional_session() as session:
            # Verify DD exists
            dd = session.query(DueDiligence).filter(DueDiligence.id == dd_id).first()
            if not dd:
                return func.HttpResponse(
                    json.dumps({"error": f"Due diligence {dd_id} not found"}),
                    status_code=404,
                    mimetype="application/json"
                )

            # Get documents to check
            # Explicitly specify the join condition since Document has multiple FKs to Folder
            docs_query = (
                session.query(Document)
                .join(Folder, Document.folder_id == Folder.id)
                .filter(
                    Folder.dd_id == dd_id,
                    Document.is_original == False
                )
            )

            if doc_ids:
                docs_query = docs_query.filter(Document.id.in_(doc_ids))

            documents = docs_query.all()

            # File types that get converted to PDF
            CONVERTIBLE_TYPES = {"pptx", "docx", "xlsx"}

            for doc in documents:
                file_type_lower = doc.type.lower()
                is_convertible = file_type_lower in CONVERTIBLE_TYPES

                # Update status to checking
                doc.readability_status = "checking"
                doc.conversion_status = "pending" if is_convertible else None
                session.commit()

                # Check readability (and convert to PDF if applicable)
                is_readable, error_msg, converted_doc_id = check_document_readability(
                    str(doc.id),
                    doc.type,
                    doc.original_file_name,
                    str(doc.folder_id),
                    session
                )

                # Update document status
                if is_readable:
                    doc.readability_status = "ready"
                    doc.readability_error = None
                    summary["ready"] += 1

                    # Handle conversion result for convertible types
                    if converted_doc_id:
                        doc.converted_doc_id = uuid.UUID(converted_doc_id)
                        doc.conversion_status = "converted"
                    elif is_convertible:
                        # Document was readable but conversion failed
                        doc.conversion_status = "failed"
                else:
                    doc.readability_status = "failed"
                    doc.readability_error = error_msg
                    summary["failed"] += 1
                    if is_convertible:
                        doc.conversion_status = "failed"

                session.commit()

                results.append({
                    "doc_id": str(doc.id),
                    "filename": doc.original_file_name,
                    "file_type": doc.type,
                    "status": doc.readability_status,
                    "error": doc.readability_error,
                    "converted_doc_id": str(doc.converted_doc_id) if doc.converted_doc_id else None,
                    "conversion_status": doc.conversion_status
                })

                logging.info(f"Readability check for {doc.original_file_name}: {doc.readability_status}")

            response = {
                "dd_id": str(dd_id),
                "total_documents": len(documents),
                "results": results,
                "summary": summary
            }

            return func.HttpResponse(
                json.dumps(response),
                status_code=200,
                mimetype="application/json"
            )

    except Exception as e:
        logging.error(f"Error in readability check: {str(e)}")
        logging.exception("Full traceback:")
        return func.HttpResponse(
            json.dumps({"error": str(e)}),
            status_code=500,
            mimetype="application/json"
        )
